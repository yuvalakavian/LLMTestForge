import os
import cv2
import pytesseract
import numpy as np
import io
from pptx import Presentation
from PIL import Image

# Set Tesseract path (Only for Windows, update if necessary)
# pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    all_text = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_text = f"--- Slide {slide_number} ---\n"
        
        # Extract text from slide shapes
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"

            # Extract text from images
            if shape.shape_type == 13:  # Shape type 13 = Picture
                image = extract_image_from_shape(shape)
                if image:
                    ocr_text = extract_text_from_image(image)
                    slide_text += f"\n[OCR from image]:\n{ocr_text}\n"

        all_text.append(slide_text)

    return "\n".join(all_text)

def extract_image_from_shape(shape):
    """Extracts an image from a PowerPoint shape."""
    try:
        if not hasattr(shape, "image"):
            return None  # Skip if shape has no image

        image_stream = io.BytesIO(shape.image.blob)
        image = Image.open(image_stream)

        return image if image else None
    except Exception as e:
        print(f"Error extracting image: {e}")
        return None

def extract_text_from_image(image):
    """Uses OCR to extract text from an image."""
    try:
        # Convert PIL image to a NumPy array
        img_array = np.array(image)

        # Ensure image is not empty
        if img_array is None or img_array.size == 0:
            return "Error: Empty image received"

        # Convert to grayscale
        gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)

        # Apply thresholding for noise reduction
        processed_image = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1]

        # Perform OCR
        return pytesseract.image_to_string(processed_image)
    except Exception as e:
        return f"Error processing image: {e}"

if __name__ == "__main__":
    pptx_file = "example.pptx"  # Change to your PowerPoint file
    extracted_text = extract_text_from_pptx(pptx_file)
    
    # Save extracted text to a file
    with open("extracted_text.txt", "w", encoding="utf-8") as f:
        f.write(extracted_text)
    
    print("Extraction complete. Text saved in 'extracted_text.txt'.")