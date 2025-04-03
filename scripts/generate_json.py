import argparse
import openai
import json
import fitz  # PyMuPDF
import io
import time
import sys
from pptx import Presentation

def read_api_key(file_path="api_key.txt"):
    with open(file_path, "r") as f:
        return f.read().strip()
    
def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    
    return "\n".join(text)

def compress_pdf_to_text(input_pdf_path):
    # Open the input PDF
    doc = fitz.open(input_pdf_path)
    full_text = ""

    # Iterate through each page and apply optimizations
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)

        # Compress images and other content
        # Remove all annotations on the page
        for annot in page.annots():
            page.delete_annot(annot)
        # page.set_(None)  # Remove annotations to save space
        page.clean_contents()  # Clean up page contents

        # Extract text from the page and append it to the string
        full_text += page.get_text("text")  # Extract text as plain text

    return full_text


def generate_content(
    generate_type, initial_prompt, response_structure, text_input
) -> int:
    # Step 1: Read API key from file
    api_key = read_api_key()

    # Initialize OpenAI Client
    openai_client = openai.OpenAI(api_key=api_key)
    
    instructions = "You are an assistant that generates tests and summary as text input in JSON format."
    
    if generate_type == "summary":
        instructions+= "The generated response must contain at least 1,000 words in JSON format."

    # Step 2: Create an Assistant
    assistant = openai_client.beta.assistants.create(
        name="Test/Summary Generator",
        instructions=instructions,
        model="o3-mini",
    )
    assistant_id = assistant.id

    print(f"Assistant created: {assistant_id}")

    # Step 3: Create a Thread
    thread = openai_client.beta.threads.create()
    thread_id = thread.id

    print(f"Thread created: {thread_id}")

    content = f"{initial_prompt} The {generate_type}: {text_input} return in JSON format acording to the format randomly: {response_structure}"
    # Step 5: Send a Message with Text Input
    message = openai_client.beta.threads.messages.create(
        thread_id=thread_id,
        role="user",
        content=content,
    )

    print("Message sent to assistant.")

    # Step 6: Run the Assistant
    run = openai_client.beta.threads.runs.create(
        thread_id=thread_id,
        assistant_id=assistant_id,
        max_completion_tokens=20000,
    )

    print("Processing...")

    # Step 7: Wait for Completion & Retrieve the Response
    while True:
        run_status = openai_client.beta.threads.runs.retrieve(
            thread_id=thread_id, run_id=run.id
        )
        if run_status.status == "completed":
            print("Processing completed.")
            break
        elif run_status.status == "failed":
            print("Error: Processing failed.")
            print(run_status)
            exit(1)
        print(run_status)
        time.sleep(2)  # Wait before checking again

    # Step 8: Fetch Messages
    messages = openai_client.beta.threads.messages.list(thread_id=thread_id)

    # Extract assistant response
    response_text = None
    for msg in messages.data:
        if msg.role == "assistant":
            for content in msg.content:
                if content.type == "text":
                    response_text = content.text.value
                    break

    if not response_text:
        print("No response received.")
        exit(1)

    # Step 9: Save Response to JSON File
    cleaned_text = (
        response_text.replace("```json\n", "").replace("\n```", "").replace("\n", "")
    )

    # Parse the cleaned text as a JSON string
    parsed_json = json.loads(cleaned_text)

    with open("output/response.json", "w", encoding="utf-8") as json_file:
        json.dump(parsed_json, json_file, indent=4, ensure_ascii=False)

    print("Response saved to response.json")

    return 0


def get_prompt(prompt_type, params, custom_prompt_arg=None):
    """Generate the appropriate prompt for OpenAI based on the request type."""
    # If a custom prompt is provided via cmd, use it; otherwise use default.
    if custom_prompt_arg:
        custom_prompt = custom_prompt_arg
    
    if prompt_type == "test":
        num_of_american = params.get("num_of_american", 8)
        num_of_open = params.get("num_of_open", 3)
        test_prompt = f"""
        Create a new and different test inspired by the material in the files,
        which will contain {num_of_american} American questions and {num_of_open} open-ended ones.
        Notes: Make sure that the answers are correct and in Hebrew. Make the questions hard. {custom_prompt}
        """
        return test_prompt
    elif prompt_type == "summary":
        # Updated prompt with explicit instructions to lengthen the summary.
        summary_prompt = f"""        
        Create a summary for me based on the material in the files.
        Make sure the summary is rich in content and well organized.
        Replace the subjects from the json with actual titles.
        Make the response as lengthy as possible, min 20 percent of actual size.
        Notes: Ensure that the summary is detailed, rich in content, and written in Hebrew. {custom_prompt}
        """
        return summary_prompt

def parse_arguments():
    """Parse command-line arguments."""
    parser = argparse.ArgumentParser(
        description="Generate test or summary from PDF or PPTX files."
    )
    parser.add_argument(
        "--generate-type", "-g",
        choices=["test", "summary"],
        required=True,
        help="Specify whether to generate a 'test' or a 'summary'."
    )
    parser.add_argument(
        "--file-type", "-f",
        choices=["pdf", "pptx"],
        required=True,
        help="Specify the file type (pdf or pptx)."
    )
    parser.add_argument(
        "--input-file", "-i",
        required=True,
        help="Path to the input PDF or PPTX file."
    )
    parser.add_argument(
        "--custom-prompt", "-c",
        help="Optional custom prompt to override the default prompt instructions."
    )
    return parser.parse_args()

if __name__ == "__main__":
    args = parse_arguments()
    generate_type = args.generate_type
    file_type = args.file_type
    input_file = args.input_file
    custom_prompt_arg = args.custom_prompt

    print(f"Generate type: {generate_type} | File type: {file_type} | Input file: {input_file}")
    if custom_prompt_arg:
        print(f"Custom prompt provided: {custom_prompt_arg}")

    # Load the appropriate response structure JSON
    response_structure_file = "test_json_structure.json" if generate_type == "test" else "summary_json_structure.json"
    with open(response_structure_file, "r", encoding="utf-8") as json_file:
        response_structure = json.load(json_file)

    # Define the initial prompt parameters
    params = {"num_of_american": 8, "num_of_open": 3} if generate_type == "test" else {}
    initial_prompt = get_prompt(prompt_type=generate_type, params=params, custom_prompt_arg=custom_prompt_arg)

    # Extract text from the input file
    if file_type == "pdf":
        total_input = compress_pdf_to_text(input_file)
    elif file_type == "pptx":
        total_input = extract_text_from_pptx(input_file)
    else:
        print("Error: Unsupported file type.")
        sys.exit(1)

    # Generate content
    result = generate_content(
        generate_type=generate_type,
        initial_prompt=initial_prompt,
        response_structure=response_structure,
        text_input=total_input,
    )

    print("Exit code:", result)
